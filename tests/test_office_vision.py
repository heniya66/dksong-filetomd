"""test_office_vision.py — 비-PDF 포맷 이미지 vision 해설 삽입(하이브리드) 검증.

검증 범위:
  (a) vision=False(기본) 시 기존 markitdown 텍스트 변환과 동일(회귀 가드).
  (b) 이미지 포함 pptx/docx/xlsx fixture 를 코드로 생성 후 vision=True 시 해설 블록
      삽입 확인(vision 엔진은 mock — 실제 LLM 호출 없음).
  (c) vision 실패 이미지 graceful skip(파이프라인 비중단).
  (d) hwpx(ZIP) 말미 모아삽입 + hwp(바이너리) 텍스트만(vision skip) 경로.

실제 LLM(Large Language Model) 호출 없음 — 모두 mock. markitdown/pptx/docx/openpyxl
미설치 시 해당 테스트는 graceful skip.

실행:
    .venv/bin/python -m pytest tests/test_office_vision.py -v
"""
from __future__ import annotations

import importlib.util
import io
import os
import sys
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if _ROOT not in sys.path:
    sys.path.insert(0, _ROOT)

import fmdw  # noqa: E402
from fmdw import markitdown_pipeline, hwp_pipeline, office_vision  # noqa: E402


def _has(mod: str) -> bool:
    return importlib.util.find_spec(mod) is not None


def _fake_vision(prompt, image_path, provider=None):
    """이미지 1장당 가짜 해설 1줄 반환(실제 LLM 호출 대체)."""
    return "FAKE 해설: 테스트 이미지"


def _failing_vision(prompt, image_path, provider=None):
    """항상 실패하는 vision 엔진(graceful skip 검증용)."""
    raise RuntimeError("simulated vision failure")


def _png_bytes(color=(0, 0, 255)) -> bytes:
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (16, 16), color).save(buf, format="PNG")
    return buf.getvalue()


def _png_file(tmp: Path, name="img.png", color=(0, 0, 255)) -> Path:
    p = tmp / name
    p.write_bytes(_png_bytes(color))
    return p


# ──────────────────────────────────────────────────────────────────────────────
# (a) vision=False 회귀 — 기존 출력과 동일
# ──────────────────────────────────────────────────────────────────────────────

@unittest.skipUnless(_has("markitdown") and _has("pptx"), "markitdown/pptx 미설치")
class TestVisionOffRegression(unittest.TestCase):

    def setUp(self):
        self._td = tempfile.TemporaryDirectory()
        self.tmp = Path(self._td.name)

    def tearDown(self):
        self._td.cleanup()

    def _make_pptx_with_image(self) -> Path:
        from pptx import Presentation
        from pptx.util import Inches

        ip = _png_file(self.tmp)
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = "SlideTextMarker"
        s.shapes.add_picture(str(ip), Inches(1), Inches(2), Inches(2), Inches(2))
        p = self.tmp / "deck.pptx"
        prs.save(p)
        return p

    def test_vision_false_equals_baseline(self):
        """vision=False(기본) 출력이 vision 인자 없는 변환과 byte-identical."""
        src = self._make_pptx_with_image()
        out_default = self.tmp / "default.md"
        out_false = self.tmp / "false.md"
        r1 = fmdw.convert_file(src, output_path=out_default)
        r2 = fmdw.convert_file(src, output_path=out_false, vision=False)
        self.assertIsNotNone(r1)
        self.assertIsNotNone(r2)
        self.assertEqual(
            out_default.read_text(encoding="utf-8"),
            out_false.read_text(encoding="utf-8"),
        )
        # 해설 블록이 들어가지 않았는지 확인.
        self.assertNotIn("vision 해설", out_false.read_text(encoding="utf-8"))

    def test_vision_false_never_calls_engine(self):
        """vision=False 면 vision 엔진(resolve)을 아예 호출하지 않는다."""
        src = self._make_pptx_with_image()
        with patch.object(office_vision, "_resolve_vision_fn") as mock_resolve:
            fmdw.convert_file(src, output_path=self.tmp / "x.md", vision=False)
            mock_resolve.assert_not_called()


# ──────────────────────────────────────────────────────────────────────────────
# (b) vision=True 해설 삽입 — pptx(정밀)/docx(말미)/xlsx(시트별)
# ──────────────────────────────────────────────────────────────────────────────

@unittest.skipUnless(_has("markitdown") and _has("pptx"), "markitdown/pptx 미설치")
class TestPptxVisionInjection(unittest.TestCase):

    def setUp(self):
        self._td = tempfile.TemporaryDirectory()
        self.tmp = Path(self._td.name)

    def tearDown(self):
        self._td.cleanup()

    def _make_pptx(self) -> Path:
        from pptx import Presentation
        from pptx.util import Inches

        ip = _png_file(self.tmp)
        prs = Presentation()
        # slide 1: text only
        s1 = prs.slides.add_slide(prs.slide_layouts[5])
        s1.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = "SLIDE1_ONLY_TEXT"
        # slide 2: text + image
        s2 = prs.slides.add_slide(prs.slide_layouts[5])
        s2.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = "SLIDE2_HAS_PIC"
        s2.shapes.add_picture(str(ip), Inches(1), Inches(2), Inches(2), Inches(2))
        p = self.tmp / "deck.pptx"
        prs.save(p)
        return p

    def test_pptx_vision_inserts_after_slide2(self):
        out = self.tmp / "deck.md"
        with patch("fmdw.ollama_extractor.extract_image", side_effect=_fake_vision):
            res = fmdw.convert_file(self.tmp_src(), output_path=out, vision=True)
        self.assertIsNotNone(res)
        text = out.read_text(encoding="utf-8")
        self.assertIn("FAKE 해설", text)
        self.assertIn("vision 해설", text)
        # 해설 블록이 Slide 2 섹션 뒤(=Slide 2 텍스트 이후)에 위치하는지: SLIDE2 마커가
        # 해설 블록보다 앞에 와야 한다.
        self.assertLess(text.index("SLIDE2_HAS_PIC"), text.index("FAKE 해설"))
        # 라벨에 Slide 2 표기.
        self.assertIn("(Slide 2)", text)

    def tmp_src(self) -> Path:
        if not hasattr(self, "_src"):
            self._src = self._make_pptx()
        return self._src

    def test_pptx_vision_uses_provider_override(self):
        """provider override 인자가 vision 엔진까지 전달되는지 확인."""
        seen = {}

        def capture(prompt, image_path, provider=None):
            seen["provider"] = provider
            return "desc"

        out = self.tmp / "deck.md"
        with patch("fmdw.ollama_extractor.extract_image", side_effect=capture):
            fmdw.convert_file(self.tmp_src(), output_path=out, vision=True, provider="gemini")
        self.assertEqual(seen.get("provider"), "gemini")


@unittest.skipUnless(_has("markitdown") and _has("docx"), "markitdown/docx 미설치")
class TestDocxVisionInjection(unittest.TestCase):

    def setUp(self):
        self._td = tempfile.TemporaryDirectory()
        self.tmp = Path(self._td.name)

    def tearDown(self):
        self._td.cleanup()

    def _make_docx(self) -> Path:
        from docx import Document
        from docx.shared import Inches

        ip = _png_file(self.tmp)
        doc = Document()
        doc.add_heading("DocxTitleMarker", level=1)
        doc.add_paragraph("DocxBodyTextMarker")
        doc.add_picture(str(ip), width=Inches(1))
        p = self.tmp / "doc.docx"
        doc.save(p)
        return p

    def test_docx_vision_appends_trailing_section(self):
        src = self._make_docx()
        out = self.tmp / "doc.md"
        with patch("fmdw.ollama_extractor.extract_image", side_effect=_fake_vision):
            res = fmdw.convert_file(src, output_path=out, vision=True)
        self.assertIsNotNone(res)
        text = out.read_text(encoding="utf-8")
        # 본문 텍스트 보존 + 말미 해설 섹션 추가.
        self.assertIn("DocxBodyTextMarker", text)
        self.assertIn("## 이미지 해설", text)
        self.assertIn("FAKE 해설", text)
        # 말미 모아삽입: 해설 섹션이 본문 텍스트 뒤.
        self.assertLess(text.index("DocxBodyTextMarker"), text.index("## 이미지 해설"))

    def test_docx_vision_failure_graceful_skip(self):
        """vision 실패 이미지는 skip — 텍스트 결과 보존, 빈 해설 섹션 미생성."""
        src = self._make_docx()
        out = self.tmp / "doc.md"
        with patch("fmdw.ollama_extractor.extract_image", side_effect=_failing_vision):
            res = fmdw.convert_file(src, output_path=out, vision=True)
        self.assertIsNotNone(res)
        text = out.read_text(encoding="utf-8")
        self.assertIn("DocxBodyTextMarker", text)
        # 전부 실패 → 빈 해설 섹션 추가 금지.
        self.assertNotIn("## 이미지 해설", text)


@unittest.skipUnless(_has("markitdown") and _has("openpyxl"), "markitdown/openpyxl 미설치")
class TestXlsxVisionInjection(unittest.TestCase):

    def setUp(self):
        self._td = tempfile.TemporaryDirectory()
        self.tmp = Path(self._td.name)

    def tearDown(self):
        self._td.cleanup()

    def _make_xlsx(self) -> Path:
        import openpyxl
        from openpyxl.drawing.image import Image as XLImage

        ip = _png_file(self.tmp)
        wb = openpyxl.Workbook()
        ws1 = wb.active
        ws1.title = "DataSheet"
        ws1["A1"] = "HEADER1"
        ws1["A2"] = "ROWVAL1"
        ws2 = wb.create_sheet("ImageSheet")
        ws2["A1"] = "HEADER2"
        img = XLImage(str(ip))
        img.anchor = "B2"
        ws2.add_image(img)
        p = self.tmp / "book.xlsx"
        wb.save(p)
        return p

    def test_xlsx_vision_inserts_in_sheet_section(self):
        src = self._make_xlsx()
        out = self.tmp / "book.md"
        with patch("fmdw.ollama_extractor.extract_image", side_effect=_fake_vision):
            res = fmdw.convert_file(src, output_path=out, vision=True)
        self.assertIsNotNone(res)
        text = out.read_text(encoding="utf-8")
        self.assertIn("ROWVAL1", text)
        self.assertIn("FAKE 해설", text)
        self.assertIn("(Sheet ImageSheet)", text)
        # 해설은 ImageSheet 섹션 뒤에 위치.
        self.assertLess(text.index("## ImageSheet"), text.index("FAKE 해설"))


# ──────────────────────────────────────────────────────────────────────────────
# (c) graceful skip — augment_with_vision 단위(엔진 무관)
# ──────────────────────────────────────────────────────────────────────────────

class TestGracefulDegrade(unittest.TestCase):

    def setUp(self):
        self._td = tempfile.TemporaryDirectory()
        self.tmp = Path(self._td.name)

    def tearDown(self):
        self._td.cleanup()

    def test_augment_no_images_returns_input_unchanged(self):
        """ZIP 에 이미지가 없으면 입력 markdown 그대로(무해)."""
        import zipfile

        z = self.tmp / "empty.docx"
        with zipfile.ZipFile(z, "w") as zf:
            zf.writestr("word/document.xml", "<xml/>")
        md = "# Title\n\nbody"
        out = office_vision.augment_with_vision(md, z, "docx", vision_fn=_fake_vision)
        self.assertEqual(out, md)

    def test_augment_non_zip_safe_degrade(self):
        """비-ZIP 파일(손상/구형 hwp) → 입력 markdown 그대로(safe degrade)."""
        f = self.tmp / "fake.hwpx"
        f.write_bytes(b"\x00not-a-zip")
        md = "본문 텍스트"
        out = office_vision.augment_with_vision(md, f, "hwpx", vision_fn=_fake_vision)
        self.assertEqual(out, md)

    def test_augment_all_images_fail_no_section(self):
        """모든 이미지 vision 실패 → 빈 해설 섹션 미생성, 입력 그대로."""
        import zipfile

        z = self.tmp / "imgs.docx"
        with zipfile.ZipFile(z, "w") as zf:
            zf.writestr("word/media/image1.png", _png_bytes())
        md = "본문"
        out = office_vision.augment_with_vision(md, z, "docx", vision_fn=_failing_vision)
        self.assertEqual(out, md)

    def test_augment_engine_load_failure_preserves_text(self):
        """vision 엔진 로드 실패 시 텍스트 결과 보존(예외 전파 금지)."""
        import zipfile

        z = self.tmp / "imgs.docx"
        with zipfile.ZipFile(z, "w") as zf:
            zf.writestr("word/media/image1.png", _png_bytes())
        md = "본문 보존"
        with patch.object(office_vision, "_resolve_vision_fn", side_effect=ImportError("no engine")):
            out = office_vision.augment_with_vision(md, z, "docx")
        self.assertEqual(out, md)

    def test_extract_zip_media_filters_non_images(self):
        """ZIP media 추출은 이미지 확장자만 통과(.xml 등 제외)."""
        import zipfile

        z = self.tmp / "m.pptx"
        with zipfile.ZipFile(z, "w") as zf:
            zf.writestr("ppt/media/image1.png", _png_bytes())
            zf.writestr("ppt/media/notes.xml", "<xml/>")
            zf.writestr("ppt/slides/slide1.xml", "<xml/>")
        media = office_vision.extract_zip_media(z, "pptx")
        names = [n for n, _ in media]
        self.assertEqual(names, ["ppt/media/image1.png"])


# ──────────────────────────────────────────────────────────────────────────────
# (d) hwpx(ZIP) 말미 모아삽입 + hwp(바이너리) 텍스트만
# ──────────────────────────────────────────────────────────────────────────────

# R11(2026-07-15): 이 클래스는 patch("pyhwp2md.convert") 를 쓰므로 pyhwp2md 모듈
# 자체가 import 가능해야 한다(선택 의존성). 미설치 venv 에서는 ModuleNotFoundError
# 로 '실패'하던 것을 정당한 환경성 skip 으로 전환 — 설치 시 자동 재활성.
@unittest.skipIf(importlib.util.find_spec("pyhwp2md") is None,
                 "pyhwp2md 미설치 — HWP 변환 선택 의존성(환경성 skip, R11)")
class TestHwpVisionPaths(unittest.TestCase):

    def setUp(self):
        self._td = tempfile.TemporaryDirectory()
        self.tmp = Path(self._td.name)

    def tearDown(self):
        self._td.cleanup()

    def _make_hwpx_with_image(self) -> Path:
        """ZIP 구조 hwpx(BinData/ 이미지 포함)를 직접 만든다(실 hwpx 생성 도구 불필요)."""
        import zipfile

        z = self.tmp / "doc.hwpx"
        with zipfile.ZipFile(z, "w") as zf:
            zf.writestr("Contents/content.hpf", "<hpf/>")
            zf.writestr("BinData/image1.png", _png_bytes())
        return z

    def test_hwpx_vision_appends_trailing(self):
        """hwpx vision=True → 말미 `## 이미지 해설` 모아삽입(convert 는 stub)."""
        src = self._make_hwpx_with_image()
        out = self.tmp / "doc.md"
        with patch("pyhwp2md.convert", return_value="HWPX_BODY 본문"), \
             patch("fmdw.ollama_extractor.extract_image", side_effect=_fake_vision):
            res = hwp_pipeline.convert_hwp(src, out, vision=True)
        self.assertIsNotNone(res)
        text = out.read_text(encoding="utf-8")
        self.assertIn("HWPX_BODY", text)
        self.assertIn("## 이미지 해설", text)
        self.assertIn("FAKE 해설", text)

    def test_hwpx_vision_false_is_baseline(self):
        """hwpx vision=False → 기존 텍스트만(해설 섹션 없음)."""
        src = self._make_hwpx_with_image()
        out = self.tmp / "doc.md"
        with patch("pyhwp2md.convert", return_value="HWPX_BODY 본문"):
            res = hwp_pipeline.convert_hwp(src, out, vision=False)
        self.assertIsNotNone(res)
        text = out.read_text(encoding="utf-8")
        self.assertIn("HWPX_BODY", text)
        self.assertNotIn("## 이미지 해설", text)

    def test_hwp_binary_vision_text_only(self):
        """hwp(바이너리) vision=True 여도 이미지 추출 skip → 텍스트만(엔진 미호출)."""
        src = self.tmp / "doc.hwp"
        src.write_bytes(b"\x00fake-hwp-binary")
        out = self.tmp / "doc.md"
        with patch("pyhwp2md.convert", return_value="HWP_BINARY_BODY"), \
             patch("fmdw.ollama_extractor.extract_image", side_effect=_fake_vision) as mock_v:
            res = hwp_pipeline.convert_hwp(src, out, vision=True)
        self.assertIsNotNone(res)
        text = out.read_text(encoding="utf-8")
        self.assertIn("HWP_BINARY_BODY", text)
        self.assertNotIn("## 이미지 해설", text)
        # hwp 바이너리는 vision 엔진을 절대 호출하지 않는다.
        mock_v.assert_not_called()


# ──────────────────────────────────────────────────────────────────────────────
# (e) convert_office vision 인자 전파(라우팅 단위)
# ──────────────────────────────────────────────────────────────────────────────

class TestVisionRouting(unittest.TestCase):

    def setUp(self):
        self._td = tempfile.TemporaryDirectory()
        self.tmp = Path(self._td.name)

    def tearDown(self):
        self._td.cleanup()

    def test_convert_file_forwards_vision_to_office(self):
        """convert_file(vision=True) 가 convert_office 까지 vision 전파."""
        src = self.tmp / "f.pptx"
        src.write_text("dummy", encoding="utf-8")
        captured = {}

        def fake_office(input_path, output_path, *, vision=False, **kw):
            captured["vision"] = vision
            Path(output_path).write_text("# ok", encoding="utf-8")
            return Path(output_path)

        with patch.object(markitdown_pipeline, "convert_office", side_effect=fake_office):
            fmdw.convert_file(src, output_dir=self.tmp, vision=True)
        self.assertTrue(captured.get("vision"))

    def test_non_vision_format_text_only(self):
        """csv 등 vision 비대상 포맷은 vision=True 여도 텍스트 그대로(_apply_office_vision noop)."""
        # _apply_office_vision 가 .csv 에 대해 markdown 을 그대로 반환하는지 직접 검증.
        out = markitdown_pipeline._apply_office_vision(
            "csv text", self.tmp / "x.csv"
        )
        self.assertEqual(out, "csv text")


if __name__ == "__main__":
    unittest.main(verbosity=2)
