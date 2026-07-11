"""test_convert_file.py — 다포맷 통합 변환 + 단일 진입점 convert_file 회귀 가드.

검증 범위:
  (1) import fmdw 가 markitdown/pyhwp2md 미설치 환경에서도 성공(지연 import).
  (2) convert_file 확장자 라우팅 + 출력 경로(output_path / output_dir / 인접) 결정.
  (3) 미지원 확장자 → ValueError(지원 목록 안내), 미존재 파일 → FileNotFoundError.
  (4) optional 의존성 미설치 시 사용 시점에 친절한 ImportError(graceful).
  (5) markitdown 설치 시 실제 office 변환(docx/pptx/xlsx/csv/html) → .md 생성.
  (6) docx engine="vision" 경로(LLM 은 stub) → output_path 저장 + truncation→.partial.md.
  (7) PDF 라우팅이 기본 프롬프트로 pdf_pipeline.convert_pdf 를 호출(실제 LLM 호출 없음 — stub).

실제 LLM 호출 없음. markitdown 은 설치되면 실제 변환(로컬, 네트워크 무관), 미설치면 graceful 테스트.

실행:
    .venv/bin/python -m pytest tests/test_convert_file.py -v
"""
from __future__ import annotations

import importlib
import importlib.util
import os
import sys
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if _ROOT not in sys.path:
    sys.path.insert(0, _ROOT)

import fmdw  # noqa: E402
from fmdw import markitdown_pipeline, hwp_pipeline  # noqa: E402


def _has_markitdown() -> bool:
    return importlib.util.find_spec("markitdown") is not None


def _has_pyhwp2md() -> bool:
    return importlib.util.find_spec("pyhwp2md") is not None


# ──────────────────────────────────────────────────────────────────────────────
# (1) import 성공 + 공개 API 표면
# ──────────────────────────────────────────────────────────────────────────────

class TestPublicApi(unittest.TestCase):

    def test_import_fmdw_succeeds_without_optional_deps(self):
        """import fmdw 시점에 markitdown/pyhwp2md 를 끌어오지 않는다(지연 import)."""
        # fmdw 는 이미 import 됨. 무거운 optional 의존성이 끌려오지 않았는지 확인.
        # (설치되어 있어도, import fmdw 만으로는 로드되지 않아야 한다.)
        # find_spec 만으로는 로드 안 되므로 sys.modules 직접 확인.
        # 단, 다른 테스트가 먼저 로드했을 수 있어 본 검증은 fresh 서브프로세스가 정확하나
        # 여기서는 export 표면만 확정 검증한다.
        for name in (
            "convert_file", "convert_pdf", "convert_docx", "convert_office",
            "convert_pptx", "convert_xlsx", "convert_html", "convert_csv",
            "convert_hwp",
        ):
            self.assertTrue(hasattr(fmdw, name), f"fmdw.{name} 누락")
            self.assertIn(name, fmdw.__all__, f"__all__ 에 {name} 누락")

    def test_supported_extensions(self):
        exts = fmdw.SUPPORTED_EXTENSIONS
        for e in (".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".html", ".htm", ".csv", ".hwp", ".hwpx"):
            self.assertIn(e, exts, f"{e} 누락")

    def test_default_pdf_prompt_has_range_placeholders(self):
        """기본 PDF 프롬프트는 청크 범위 치환용 {start}/{end} 를 포함해야 한다."""
        self.assertIn("{start}", fmdw.DEFAULT_PDF_PROMPT)
        self.assertIn("{end}", fmdw.DEFAULT_PDF_PROMPT)

    def test_import_lazy_in_subprocess(self):
        """fresh 프로세스에서 import fmdw 후 markitdown/pyhwp2md 미로드 확인."""
        code = (
            "import sys; import fmdw; "
            "print('markitdown' in sys.modules, 'pyhwp2md' in sys.modules)"
        )
        import subprocess
        out = subprocess.run(
            [sys.executable, "-c", code],
            capture_output=True, text=True, cwd=_ROOT,
        )
        self.assertEqual(out.returncode, 0, out.stderr)
        self.assertIn("False False", out.stdout)


# ──────────────────────────────────────────────────────────────────────────────
# (2)(3) 라우팅 / 출력 경로 / 에러
# ──────────────────────────────────────────────────────────────────────────────

class TestRoutingAndOutputPath(unittest.TestCase):

    def setUp(self):
        self._td = tempfile.TemporaryDirectory()
        self.tmp = Path(self._td.name)

    def tearDown(self):
        self._td.cleanup()

    def _touch(self, name: str) -> Path:
        p = self.tmp / name
        p.write_text("dummy", encoding="utf-8")
        return p

    def test_unsupported_extension_raises_valueerror(self):
        src = self._touch("a.txt")
        with self.assertRaises(ValueError) as ctx:
            fmdw.convert_file(src, output_dir=self.tmp)
        # 지원 목록 안내 포함.
        self.assertIn(".pdf", str(ctx.exception))

    def test_missing_file_raises(self):
        with self.assertRaises(FileNotFoundError):
            fmdw.convert_file(self.tmp / "nope.docx", output_dir=self.tmp)

    def test_output_dir_routing(self):
        """output_dir 지정 시 <output_dir>/<stem>.md 로 라우팅 (변환 함수는 stub)."""
        src = self._touch("report.pptx")
        outdir = self.tmp / "02_processed_md"
        captured = {}

        def fake_office(input_path, output_path, **kw):
            captured["out"] = Path(output_path)
            Path(output_path).parent.mkdir(parents=True, exist_ok=True)
            Path(output_path).write_text("# ok", encoding="utf-8")
            return Path(output_path)

        with patch.object(markitdown_pipeline, "convert_office", side_effect=fake_office):
            res = fmdw.convert_file(src, output_dir=outdir)
        self.assertEqual(captured["out"], outdir / "report.md")
        self.assertEqual(res, outdir / "report.md")
        self.assertTrue((outdir / "report.md").exists())

    def test_output_path_overrides_output_dir(self):
        src = self._touch("data.csv")
        explicit = self.tmp / "custom" / "x.md"
        captured = {}

        def fake_office(input_path, output_path, **kw):
            captured["out"] = Path(output_path)
            return Path(output_path)

        with patch.object(markitdown_pipeline, "convert_office", side_effect=fake_office):
            fmdw.convert_file(src, output_dir=self.tmp / "ignored", output_path=explicit)
        self.assertEqual(captured["out"], explicit)

    def test_no_output_specified_uses_src_sibling(self):
        src = self._touch("sheet.xlsx")
        captured = {}

        def fake_office(input_path, output_path, **kw):
            captured["out"] = Path(output_path)
            return Path(output_path)

        with patch.object(markitdown_pipeline, "convert_office", side_effect=fake_office):
            fmdw.convert_file(src)
        self.assertEqual(captured["out"], self.tmp / "sheet.md")

    def test_extension_case_insensitive(self):
        src = self._touch("DOC.DOCX")
        captured = {}

        def fake_docx(input_path, output_path, **kw):
            captured["out"] = Path(output_path)
            return Path(output_path)

        with patch.object(markitdown_pipeline, "convert_docx", side_effect=fake_docx):
            fmdw.convert_file(src, output_dir=self.tmp)
        self.assertEqual(captured["out"], self.tmp / "DOC.md")

    def test_office_extensions_route_to_office(self):
        for ext in (".pptx", ".xlsx", ".xls", ".html", ".htm", ".csv"):
            src = self._touch(f"f{ext}")
            with patch.object(markitdown_pipeline, "convert_office", return_value=Path("x")) as m:
                fmdw.convert_file(src, output_dir=self.tmp)
            self.assertTrue(m.called, f"{ext} → convert_office 미라우팅")

    def test_hwp_extensions_route_to_hwp(self):
        for ext in (".hwp", ".hwpx"):
            src = self._touch(f"f{ext}")
            with patch.object(hwp_pipeline, "convert_hwp", return_value=Path("x")) as m:
                fmdw.convert_file(src, output_dir=self.tmp)
            self.assertTrue(m.called, f"{ext} → convert_hwp 미라우팅")

    def test_pdf_routes_with_default_prompt(self):
        """PDF 라우팅 시 prompt 미지정이면 DEFAULT_PDF_PROMPT 가 convert_pdf 로 전달."""
        src = self._touch("doc.pdf")
        captured = {}

        def fake_pdf(pdf_path, prompt_template, *, output_path, **kw):
            captured["prompt"] = prompt_template
            captured["out"] = Path(output_path)
            return Path(output_path)

        # __init__ 의 convert_pdf 는 pdf_pipeline.convert_pdf 를 lazy import → 그 대상 patch.
        with patch("fmdw.pdf_pipeline.convert_pdf", side_effect=fake_pdf):
            fmdw.convert_file(src, output_dir=self.tmp)
        self.assertEqual(captured["prompt"], fmdw.DEFAULT_PDF_PROMPT)
        self.assertEqual(captured["out"], self.tmp / "doc.md")

    def test_pdf_prompt_override(self):
        src = self._touch("doc.pdf")
        captured = {}

        def fake_pdf(pdf_path, prompt_template, *, output_path, **kw):
            captured["prompt"] = prompt_template
            return Path(output_path)

        with patch("fmdw.pdf_pipeline.convert_pdf", side_effect=fake_pdf):
            fmdw.convert_file(src, output_dir=self.tmp, prompt="MY PROMPT")
        self.assertEqual(captured["prompt"], "MY PROMPT")

    def test_pdf_kwargs_forwarded(self):
        """convert_file kwargs(chunk_size 등)가 convert_pdf 로 전달."""
        src = self._touch("doc.pdf")
        captured = {}

        def fake_pdf(pdf_path, prompt_template, *, output_path, **kw):
            captured.update(kw)
            return Path(output_path)

        with patch("fmdw.pdf_pipeline.convert_pdf", side_effect=fake_pdf):
            fmdw.convert_file(src, output_dir=self.tmp, chunk_size=7, on_failure="partial")
        self.assertEqual(captured.get("chunk_size"), 7)
        self.assertEqual(captured.get("on_failure"), "partial")


# ──────────────────────────────────────────────────────────────────────────────
# (4) optional 의존성 미설치 graceful ImportError
# ──────────────────────────────────────────────────────────────────────────────

class TestGracefulImportError(unittest.TestCase):

    def setUp(self):
        self._td = tempfile.TemporaryDirectory()
        self.tmp = Path(self._td.name)

    def tearDown(self):
        self._td.cleanup()

    def test_office_missing_raises_friendly(self):
        """markitdown 미설치를 시뮬레이트 → 친절한 설치 안내 ImportError."""
        src = self.tmp / "a.pptx"
        src.write_text("dummy", encoding="utf-8")
        # _require_markitdown 이 실제 import 하므로, import 단을 직접 막는다.
        orig_import = __import__

        def blocked_import(name, *a, **k):
            if name == "markitdown" or name.startswith("markitdown."):
                raise ImportError("simulated missing markitdown")
            return orig_import(name, *a, **k)

        with patch("builtins.__import__", side_effect=blocked_import):
            with self.assertRaises(ImportError) as ctx:
                markitdown_pipeline.convert_office(src, self.tmp / "a.md")
        self.assertIn("filestomdwgem[office]", str(ctx.exception))

    def test_hwp_missing_raises_friendly(self):
        src = self.tmp / "a.hwp"
        src.write_text("dummy", encoding="utf-8")
        orig_import = __import__

        def blocked_import(name, *a, **k):
            if name == "pyhwp2md" or name.startswith("pyhwp2md."):
                raise ImportError("simulated missing pyhwp2md")
            return orig_import(name, *a, **k)

        with patch("builtins.__import__", side_effect=blocked_import):
            with self.assertRaises(ImportError) as ctx:
                hwp_pipeline.convert_hwp(src, self.tmp / "a.md")
        self.assertIn("filestomdwgem[hwp]", str(ctx.exception))

    def test_unknown_docx_engine_raises(self):
        src = self.tmp / "a.docx"
        src.write_text("dummy", encoding="utf-8")
        with self.assertRaises(ValueError):
            markitdown_pipeline.convert_docx(src, self.tmp / "a.md", engine="bogus")


# ──────────────────────────────────────────────────────────────────────────────
# (5) markitdown 실제 변환 (설치 시에만)
# ──────────────────────────────────────────────────────────────────────────────

@unittest.skipUnless(_has_markitdown(), "markitdown 미설치 — 실동작 변환 테스트 skip")
class TestRealOfficeConversion(unittest.TestCase):

    def setUp(self):
        self._td = tempfile.TemporaryDirectory()
        self.tmp = Path(self._td.name)
        self.outdir = self.tmp / "02_processed_md"

    def tearDown(self):
        self._td.cleanup()

    def test_docx_markitdown(self):
        from docx import Document
        doc = Document()
        doc.add_heading("Title One", level=1)
        doc.add_paragraph("Hello body text.")
        t = doc.add_table(rows=2, cols=2)
        t.cell(0, 0).text = "Name"
        t.cell(0, 1).text = "Value"
        t.cell(1, 0).text = "R1"
        t.cell(1, 1).text = "10K"
        src = self.tmp / "sample.docx"
        doc.save(src)

        res = fmdw.convert_file(src, output_dir=self.outdir)
        self.assertIsNotNone(res)
        self.assertEqual(res, self.outdir / "sample.md")
        text = res.read_text(encoding="utf-8")
        self.assertIn("Title One", text)
        self.assertIn("Hello body text", text)

    def test_pptx_conversion(self):
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            self.skipTest("python-pptx 미설치")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tb.text_frame.text = "Slide marker XYZ"
        src = self.tmp / "deck.pptx"
        prs.save(src)

        res = fmdw.convert_file(src, output_dir=self.outdir)
        self.assertIsNotNone(res)
        self.assertEqual(res, self.outdir / "deck.md")
        self.assertIn("Slide marker XYZ", res.read_text(encoding="utf-8"))

    def test_xlsx_conversion(self):
        try:
            import openpyxl
        except ImportError:
            self.skipTest("openpyxl 미설치")
        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = "Header"
        ws["B1"] = "Num"
        ws["A2"] = "RowVal"
        ws["B2"] = 42
        src = self.tmp / "book.xlsx"
        wb.save(src)

        res = fmdw.convert_file(src, output_dir=self.outdir)
        self.assertIsNotNone(res)
        self.assertEqual(res, self.outdir / "book.md")
        self.assertIn("RowVal", res.read_text(encoding="utf-8"))

    def test_csv_conversion(self):
        src = self.tmp / "rows.csv"
        src.write_text("col1,col2\nfoo,123\nbar,456\n", encoding="utf-8")
        res = fmdw.convert_file(src, output_dir=self.outdir)
        self.assertIsNotNone(res)
        text = res.read_text(encoding="utf-8")
        self.assertIn("foo", text)
        self.assertIn("456", text)

    def test_html_conversion(self):
        src = self.tmp / "page.html"
        src.write_text(
            "<html><body><h1>HeadingH</h1><p>ParaP content</p></body></html>",
            encoding="utf-8",
        )
        res = fmdw.convert_file(src, output_dir=self.outdir)
        self.assertIsNotNone(res)
        text = res.read_text(encoding="utf-8")
        self.assertIn("HeadingH", text)
        self.assertIn("ParaP content", text)


# ──────────────────────────────────────────────────────────────────────────────
# (6) docx engine="vision" 경로 (LLM stub)
# ──────────────────────────────────────────────────────────────────────────────

class TestDocxVisionEngine(unittest.TestCase):

    def setUp(self):
        self._td = tempfile.TemporaryDirectory()
        self.tmp = Path(self._td.name)

    def tearDown(self):
        self._td.cleanup()

    def _make_docx(self) -> Path:
        from docx import Document
        doc = Document()
        doc.add_paragraph("Vision body content")
        t = doc.add_table(rows=1, cols=2)
        t.cell(0, 0).text = "k"
        t.cell(0, 1).text = "v"
        p = self.tmp / "vdoc.docx"
        doc.save(p)
        return p

    def test_vision_clean_saves_to_output_path(self):
        import extract_docx_multimodal as dm
        import fmdw.ollama_extractor as ox
        src = self._make_docx()
        out = self.tmp / "out" / "vdoc.md"
        with patch.object(ox, "extract_text_prompt", return_value="## clean vision md"), \
             patch.object(ox, "provider_label", lambda: "stub"):
            res = fmdw.convert_file(src, output_path=out, engine="vision")
        self.assertEqual(res, out)
        self.assertTrue(out.exists())
        self.assertIn("clean vision md", out.read_text(encoding="utf-8"))
        # extract_docx_multimodal 모듈은 수정 없이 재활용됨(공개 함수 존재 확인).
        self.assertTrue(hasattr(dm, "extract_blocks_from_docx"))
        self.assertTrue(hasattr(dm, "chunk_blocks"))

    def test_vision_truncation_saves_partial(self):
        import fmdw.ollama_extractor as ox
        src = self._make_docx()
        out = self.tmp / "vdoc.md"
        trunc = "## partial\n<!-- TRUNCATED: finish_reason=length -->"
        with patch.object(ox, "extract_text_prompt", return_value=trunc), \
             patch.object(ox, "provider_label", lambda: "stub"):
            res = fmdw.convert_file(src, output_path=out, engine="vision")
        # truncation → .partial.md (완성본 위장 금지).
        self.assertEqual(res, self.tmp / "vdoc.partial.md")
        self.assertTrue((self.tmp / "vdoc.partial.md").exists())
        self.assertFalse(out.exists())


# ──────────────────────────────────────────────────────────────────────────────
# (7) hwp 실동작 (설치 시에만) — 미설치면 graceful 테스트가 (4)에서 커버
# ──────────────────────────────────────────────────────────────────────────────

@unittest.skipUnless(_has_pyhwp2md(), "pyhwp2md 미설치 — hwp 실동작 테스트 skip")
class TestRealHwpConversion(unittest.TestCase):

    def test_convert_hwp_stubbed_convert(self):
        """pyhwp2md.convert 를 stub 으로 교체해 저장 계약만 검증(실 hwp 불필요)."""
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            src = tmp / "doc.hwp"
            src.write_bytes(b"\x00fake-hwp")
            out = tmp / "doc.md"
            with patch("pyhwp2md.convert", return_value="# hwp markdown"):
                res = hwp_pipeline.convert_hwp(src, out)
            self.assertEqual(res, out)
            self.assertIn("hwp markdown", out.read_text(encoding="utf-8"))

    def test_hwpx_real_roundtrip(self):
        """python-hwpx 로 실제 .hwpx 생성 → convert_file → .md 본문 검증(실동작)."""
        if importlib.util.find_spec("hwpx") is None:
            self.skipTest("python-hwpx 미설치 — 실 hwpx 생성 불가")
        from hwpx import DocumentPlan, DocumentBlock, create_document_from_plan
        plan = DocumentPlan(
            title="HWPX Title Marker",
            blocks=[DocumentBlock(type="paragraph", data={"text": "Hello HWPX body XYZ"})],
        )
        doc = create_document_from_plan(plan)
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            src = tmp / "sample.hwpx"
            # save_to_path 우선(신 API), 없으면 save fallback(구버전 호환).
            if hasattr(doc, "save_to_path"):
                doc.save_to_path(str(src))
            else:  # pragma: no cover — 구버전 python-hwpx 호환.
                doc.save(str(src))
            outdir = tmp / "02_processed_md"
            res = fmdw.convert_file(src, output_dir=outdir)
            self.assertIsNotNone(res)
            self.assertEqual(res, outdir / "sample.md")
            txt = res.read_text(encoding="utf-8")
            self.assertIn("HWPX Title Marker", txt)
            self.assertIn("Hello HWPX body XYZ", txt)


if __name__ == "__main__":
    unittest.main(verbosity=2)
