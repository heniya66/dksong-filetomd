"""filestomdwgem 비-PDF 포맷 이미지 vision 해설 삽입 (하이브리드).

markitdown(Markit Down) 의 텍스트/표 변환 결과에 더해, 문서 안에 들어 있는
**이미지(그림/사진/도면)** 를 fmdw 의 기존 vision 엔진(PDF 와 동일 provider 추상화)으로
해설(Markdown 텍스트)해서 본문 해당 위치에 삽입한다.

대상 포맷: docx / pptx / xlsx / hwpx (ZIP(Zip archive) 구조). hwp(바이너리)는 이미지
추출 난이도가 높아 이번 단계에서는 텍스트만(vision skip — convert_hwp 가 그대로 처리).

설계 원칙 (CLAUDE.md 글로벌 룰 준수):
  - **opt-in**: 모든 진입점에서 vision=False 가 기본. False 면 기존 markitdown 텍스트
    변환과 100% 동일(이 모듈은 호출조차 되지 않는다). PDF 파이프라인은 절대 건드리지 않는다.
  - **vision 엔진 재사용**: 새 LLM(Large Language Model) 호출 코드를 만들지 않고
    fmdw.ollama_extractor.extract_image(prompt, image_path, provider) 를 그대로 사용한다
    (PDF 의 figure 추출이 쓰는 것과 동일한 provider 추상화 — ollama_cloud 기본, gemini
    fallback). provider override 인자를 그대로 전달한다.
  - **safe degrade**: 개별 이미지 vision 실패는 경고 후 skip — 파이프라인을 절대 중단하지
    않는다. vision=True 인데 이미지가 0개면 텍스트 결과를 그대로 둔다(무해).
  - **lazy import**: markitdown / python-pptx / python-docx / openpyxl / PIL 등 무거운
    optional 의존성은 함수 호출 시점에만 import 한다. `import fmdw.office_vision` 자체는
    의존성 미설치 환경에서도 성공해야 한다.

위치 매핑 (포맷별):
  - pptx(정밀): python-pptx 로 슬라이드별 그림 추출 → 해당 `<!-- Slide number: N -->`
    섹션 뒤에 그 슬라이드 이미지 해설 삽입.
  - xlsx(시트별): openpyxl 로 시트별 이미지(anchor) → 해당 `## <시트명>` 섹션 뒤에 삽입.
  - docx(근사): markitdown 이 인라인 이미지를 본문에 남기지 않으므로 정밀 인라인 매핑이
    불가능 → ZIP `word/media/` 의 이미지를 본문 끝 `## 이미지 해설` 섹션에 출현(파일명)
    순서대로 모아 삽입한다(택한 방식: append-at-end). 한계 명시.
  - hwpx(말미): ZIP media 추출 → 위치 정보를 본문 좌표로 매핑하기 어려워 본문 끝
    `## 이미지 해설` 섹션에 모아 삽입(한계 명시).
"""

from __future__ import annotations

import io
import logging
import os
import zipfile
from pathlib import Path
from typing import Callable, Optional

_log = logging.getLogger(__name__)


def _describe_lang() -> str:
    """FMDW_DESCRIBE_LANG: office 이미지 vision 해설 출력 언어(기본 "en").

    fmdw.figure_extractor._describe_lang() 과 동일한 계약(기본 "en"=영어 그대로 최종
    산출물, "ko"=한국어). 본 모듈은 무거운 optional 의존성을 지연 import 하는 설계
    원칙(상단 docstring)을 지키기 위해 figure_extractor 를 top-level import 하지 않고
    동일 로직을 그대로 복제한다(단일 env 계약은 공유, 모듈 결합은 0).
    """
    return os.getenv("FMDW_DESCRIBE_LANG", "en").strip().lower()

# 해설 블록 헤더 텍스트(말미 모아삽입 섹션 제목).
_TRAILING_SECTION_TITLE = "## 이미지 해설"

# 단일 이미지 vision 해설 프롬프트(범용 — PDF DEFAULT_PDF_PROMPT 의 이미지 규칙과 동일 취지).
# 출력 언어는 FMDW_DESCRIBE_LANG(기본 "en")로 선택한다(_describe_lang() 참조): 기본값이면
# 아래 _IMAGE_VISION_PROMPT_EN 을, "ko" 면 기존 _IMAGE_VISION_PROMPT_KO 를 사용한다.
_IMAGE_VISION_PROMPT_KO = """이 이미지를 상세하고 서술적인 한국어 Markdown 으로 설명하라.

규칙:
- 맨 앞에 '한 줄 요약'을 한 문장으로 적어라: 이 이미지가 무엇인지 전문지식 없는 사람도 알 수 있는 평범한 말로(예: '이것은 ~를 보여주는 블록 다이어그램이다').
- 단순 나열이 아니라 서술적으로 풀어라: (1) 무엇을 나타내는 이미지인지(도면/그림/표/사진 등 종류·목적), (2) 주요 구성요소·라벨 각각의 이름과 역할·기능, (3) 요소 간 연결·신호 흐름·배치 관계, (4) 전달하는 기능적 의미·맥락을 문단으로 자세히 기술.
- 비전문가도 끝까지 이해하도록 쉽고 자세히 설명하라: 구성요소·연결·흐름·각 부분의 역할을 일상 언어로 이야기하듯 풀어라. 단, 아래 정밀 전사(텍스트/숫자/표)를 줄이지 말고 그 위에 '쉬운 설명'을 더한다(정밀 데이터 누락 금지).
- 전문용어·약어는 처음 나올 때 바로 옆 괄호로 쉬운 뜻을 덧붙여라(예: 'GPIO(범용 입출력 핀)', 'net(부품을 잇는 전기 배선)').
- 표는 GFM 표로 정확히. 이미지 내부 텍스트/숫자는 OCR(Optical Character Recognition) 수준으로 전사.
- 식별번호([Fig. N], [도 N], [그림 N] 등)가 보이면 원문 표기 그대로 유지.
- 추정/환각 금지 — 쉽게 풀어 설명하되 '보이는 사실'만 적는다. 보이지 않거나 불확실하면 지어내지 말고 '불명확/판독 불가'로 명시. 판독 불가 셀은 [unreadable]. (쉬운 설명 ↔ 정확성 충돌 시 정확성 우선.)
- 코드펜스(```)로 감싸지 말고 본문만 출력.
"""

#: FMDW_DESCRIBE_LANG=en(기본)일 때 쓰는 영어 해설 프롬프트. _IMAGE_VISION_PROMPT_KO 와
#: 규칙 순서·anti-fabrication 제약은 동일하고 출력 언어만 영어.
_IMAGE_VISION_PROMPT_EN = """Describe this image in detailed, narrative English Markdown.

Rules:
- Start with a 'one-line summary': one plain sentence describing what this image is, understandable even to someone with no technical background (e.g. "This is a block diagram showing ~").
- Write narratively, not as a bare list: describe in paragraphs (1) what kind of image this is and its purpose (diagram/figure/table/photo, etc.), (2) the name and role/function of each major component/label, (3) the connections/signal flow/layout relationships between elements, and (4) the functional meaning/context it conveys.
- Explain in plain, detailed language so a non-expert can follow all the way through: describe the components/connections/flow/role of each part in everyday, narrative language. However, do not shorten the precise transcription (text/numbers/tables) below — add this easy explanation on top of it (no loss of precise data).
- The first time a technical term/abbreviation appears, add its plain meaning in parentheses right after it (e.g. 'GPIO (general-purpose input/output pin)', 'net (electrical wiring connecting components)').
- Transcribe tables precisely as GFM tables. Transcribe text/numbers inside the image at OCR(Optical Character Recognition) quality.
- If an identifier ([Fig. N], [Fig N], etc.) is visible, keep it verbatim as written in the original.
- Do not guess or fabricate — explain in plain terms but write only what is actually visible. If something is not visible or unclear, do not invent it — state clearly 'unclear/unreadable'. Mark unreadable cells as [unreadable]. (If the plain explanation conflicts with accuracy, accuracy always takes priority.)
- Do not wrap the output in code fences (```); output only the body text.
"""

# ZIP 내부 media 경로 prefix (포맷별). 소문자 비교.
_ZIP_MEDIA_PREFIXES = {
    "docx": ("word/media/",),
    "pptx": ("ppt/media/",),
    "xlsx": ("xl/media/",),
    # hwpx 구조는 배포 도구에 따라 BinData/ 또는 Contents/ 하위에 바이너리가 들어간다.
    # 이미지 확장자 필터와 병행하여 양쪽을 모두 후보로 본다.
    "hwpx": ("bindata/", "contents/"),
}

# 이미지로 간주할 확장자(소문자).
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".webp", ".emf", ".wmf"}


# ──────────────────────────────────────────────────────────────────────────────
# vision 엔진 호출 헬퍼 (fmdw.ollama_extractor.extract_image 재사용)
# ──────────────────────────────────────────────────────────────────────────────

def _resolve_vision_fn() -> Callable[..., str]:
    """이미지 → 해설 텍스트 vision 함수를 lazy import 로 가져온다.

    PDF 와 동일한 provider 추상화(ollama_cloud 기본, gemini fallback)를 쓰는
    fmdw.ollama_extractor.extract_image(prompt, image_path, provider=...) 를 반환한다.
    새 LLM 호출 코드를 만들지 않고 기존 엔진을 그대로 재사용한다.
    """
    try:
        from fmdw.ollama_extractor import extract_image  # 패키지 경로
    except ImportError:  # pragma: no cover — 직접 실행/경로 차이 대비
        from ollama_extractor import extract_image  # type: ignore
    return extract_image


def _describe_image_bytes(
    data: bytes,
    suffix: str,
    *,
    vision_fn: Callable[..., str],
    provider: Optional[str],
    prompt: str,
) -> Optional[str]:
    """이미지 바이트를 임시 파일로 써서 vision 엔진으로 해설 텍스트를 얻는다.

    extract_image 는 파일 경로를 받으므로(내부에서 base64 인코딩) 격리 임시 파일에
    바이트를 쓰고 경로를 넘긴다. 실패는 None 반환(호출자가 skip + 경고).

    Args:
        data: 이미지 raw 바이트.
        suffix: 임시 파일 확장자(.png/.jpg 등 — MIME 추정에 사용).
        vision_fn: extract_image 콜러블(테스트에서 직접 주입/mock 가능).
        provider: provider override(None 이면 모듈 기본 EXTRACT_PROVIDER).
        prompt: vision 프롬프트.

    Returns:
        해설 Markdown 텍스트(성공) 또는 None(실패/빈 응답).
    """
    import tempfile

    safe_suffix = suffix if suffix.lower() in _IMAGE_EXTS else ".png"
    tmp_path: Optional[Path] = None
    try:
        fd, name = tempfile.mkstemp(prefix="officevis_", suffix=safe_suffix)
        import os as _os

        _os.close(fd)
        tmp_path = Path(name)
        tmp_path.write_bytes(data)
        text = vision_fn(prompt, tmp_path, provider=provider)
        if not text or not str(text).strip():
            return None
        return str(text).strip()
    except Exception as e:  # noqa: BLE001 — 개별 이미지 실패는 파이프라인 비중단(safe degrade).
        _log.warning("이미지 vision 해설 실패 → skip: %s", e)
        return None
    finally:
        if tmp_path is not None:
            try:
                tmp_path.unlink(missing_ok=True)
            except Exception:  # noqa: BLE001
                pass


def _format_block(index: int, label: str, description: str) -> str:
    """이미지 해설 블록 Markdown 을 통일 형식으로 만든다.

    예:
        > **[이미지 1 · vision 해설]** (Slide 2)
        > <해설 첫 줄>
        > <해설 둘째 줄>
    """
    header = f"> **[이미지 {index} · vision 해설]**"
    if label:
        header += f" ({label})"
    body_lines = [f"> {ln}" if ln else ">" for ln in description.splitlines()]
    return header + "\n" + "\n".join(body_lines)


# ──────────────────────────────────────────────────────────────────────────────
# ZIP 기반 미디어 추출 (모든 ZIP 포맷 공통 — docx/pptx/xlsx/hwpx)
# ──────────────────────────────────────────────────────────────────────────────

def extract_zip_media(path: Path, fmt: str) -> list[tuple[str, bytes]]:
    """ZIP 구조 office/hwpx 파일에서 내부 이미지 미디어를 (이름, 바이트)로 추출한다.

    포맷별 media prefix(_ZIP_MEDIA_PREFIXES)와 이미지 확장자 필터로 후보를 거른다.
    ZIP 이 아니거나(예: 구형 .hwp 바이너리) 열기 실패 시 빈 리스트(safe degrade).

    Returns:
        [(arcname, data), ...] — ZIP 내부 정렬 순서(파일명 정렬)로 반환.
    """
    prefixes = _ZIP_MEDIA_PREFIXES.get(fmt, ())
    out: list[tuple[str, bytes]] = []
    try:
        with zipfile.ZipFile(path) as zf:
            names = sorted(zf.namelist())
            for name in names:
                low = name.lower()
                if prefixes and not any(low.startswith(p) for p in prefixes):
                    continue
                ext = Path(low).suffix
                if ext not in _IMAGE_EXTS:
                    continue
                try:
                    data = zf.read(name)
                except Exception as e:  # noqa: BLE001
                    _log.warning("ZIP 내부 이미지 읽기 실패 %s: %s", name, e)
                    continue
                if data:
                    out.append((name, data))
    except (zipfile.BadZipFile, OSError) as e:
        _log.warning("%s ZIP 미디어 추출 실패(비-ZIP/손상 의심) → 이미지 skip: %s", fmt, e)
    return out


# ──────────────────────────────────────────────────────────────────────────────
# pptx — 슬라이드별 정밀 매핑
# ──────────────────────────────────────────────────────────────────────────────

def _collect_pptx_images_by_slide(path: Path) -> dict[int, list[tuple[str, bytes]]]:
    """python-pptx 로 슬라이드(1-based)별 그림 이미지를 추출한다.

    실패/미설치 시 빈 dict(호출자가 ZIP fallback 또는 skip).
    """
    try:
        from pptx import Presentation  # lazy
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    except Exception as e:  # noqa: BLE001
        _log.warning("python-pptx 미설치/로드 실패 → pptx 이미지 skip: %s", e)
        return {}

    by_slide: dict[int, list[tuple[str, bytes]]] = {}
    try:
        prs = Presentation(str(path))
    except Exception as e:  # noqa: BLE001
        _log.warning("pptx 열기 실패 → 이미지 skip: %s", e)
        return {}

    for sidx, slide in enumerate(prs.slides, start=1):
        imgs: list[tuple[str, bytes]] = []
        for shape in slide.shapes:
            try:
                is_pic = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            except Exception:  # noqa: BLE001
                is_pic = getattr(shape, "shape_type", None) == 13
            if not is_pic:
                continue
            try:
                blob = shape.image.blob
                ext = "." + (shape.image.ext or "png").lstrip(".")
            except Exception as e:  # noqa: BLE001
                _log.warning("pptx slide %d 이미지 blob 추출 실패 → skip: %s", sidx, e)
                continue
            if blob:
                imgs.append((ext, blob))
        if imgs:
            by_slide[sidx] = imgs
    return by_slide


def _inject_pptx(
    markdown: str,
    path: Path,
    *,
    vision_fn: Callable[..., str],
    provider: Optional[str],
    prompt: str,
) -> str:
    """pptx markdown 의 각 `<!-- Slide number: N -->` 섹션 뒤에 그 슬라이드 이미지 해설 삽입.

    슬라이드 마커를 찾지 못하면(markitdown 버전 차이 등) 말미 모아삽입으로 fallback.
    """
    by_slide = _collect_pptx_images_by_slide(path)
    if not by_slide:
        return markdown

    marker = "<!-- Slide number:"
    if marker not in markdown:
        # 슬라이드 마커 없음 → 말미 모아삽입 fallback(순서 보존).
        flat: list[tuple[str, bytes, str]] = []
        for sidx in sorted(by_slide):
            for ext, blob in by_slide[sidx]:
                flat.append((ext, blob, f"Slide {sidx}"))
        return _append_trailing(markdown, flat, vision_fn=vision_fn, provider=provider, prompt=prompt)

    # 슬라이드 마커 기준으로 세그먼트 분할 후, 각 세그먼트 끝에 해당 슬라이드 해설 삽입.
    lines = markdown.splitlines()
    # 각 마커 라인의 (line_index, slide_number) 수집.
    seg_starts: list[tuple[int, int]] = []
    for i, ln in enumerate(lines):
        s = ln.strip()
        if s.startswith(marker):
            num = _parse_slide_number(s)
            if num is not None:
                seg_starts.append((i, num))

    if not seg_starts:
        flat = []
        for sidx in sorted(by_slide):
            for ext, blob in by_slide[sidx]:
                flat.append((ext, blob, f"Slide {sidx}"))
        return _append_trailing(markdown, flat, vision_fn=vision_fn, provider=provider, prompt=prompt)

    # 세그먼트 끝(다음 마커 직전, 마지막은 EOF)에 해설을 끼워넣는다. 뒤에서 앞으로 삽입해
    # 인덱스 변동을 피한다.
    counter = [0]  # 전체 이미지 번호(문서 단위 증가).
    # (insert_position, block_text) 목록 — insert_position 은 "그 줄 앞에 삽입".
    inserts: list[tuple[int, str]] = []
    for seg_i, (line_idx, slide_num) in enumerate(seg_starts):
        end = seg_starts[seg_i + 1][0] if seg_i + 1 < len(seg_starts) else len(lines)
        imgs = by_slide.get(slide_num, [])
        if not imgs:
            continue
        blocks = _describe_many(
            imgs, label=f"Slide {slide_num}", counter=counter,
            vision_fn=vision_fn, provider=provider, prompt=prompt,
        )
        if blocks:
            inserts.append((end, "\n\n" + "\n\n".join(blocks)))

    if not inserts:
        return markdown
    for pos, text in sorted(inserts, key=lambda t: t[0], reverse=True):
        lines.insert(pos, text)
    return "\n".join(lines)


def _parse_slide_number(marker_line: str) -> Optional[int]:
    """`<!-- Slide number: 3 -->` 에서 정수 3 을 파싱. 실패 시 None."""
    import re

    m = re.search(r"Slide number:\s*(\d+)", marker_line)
    if not m:
        return None
    try:
        return int(m.group(1))
    except ValueError:  # pragma: no cover
        return None


# ──────────────────────────────────────────────────────────────────────────────
# xlsx — 시트별 매핑
# ──────────────────────────────────────────────────────────────────────────────

def _collect_xlsx_images_by_sheet(path: Path) -> dict[str, list[tuple[str, bytes]]]:
    """openpyxl 로 시트명별 이미지(anchor)를 추출한다. 실패/미설치 시 빈 dict."""
    try:
        import openpyxl  # lazy
    except Exception as e:  # noqa: BLE001
        _log.warning("openpyxl 미설치/로드 실패 → xlsx 이미지 skip: %s", e)
        return {}

    by_sheet: dict[str, list[tuple[str, bytes]]] = {}
    try:
        wb = openpyxl.load_workbook(str(path))
    except Exception as e:  # noqa: BLE001
        _log.warning("xlsx 열기 실패 → 이미지 skip: %s", e)
        return {}

    for ws in wb.worksheets:
        imgs: list[tuple[str, bytes]] = []
        for im in getattr(ws, "_images", []) or []:
            blob = _xlsx_image_bytes(im)
            if not blob:
                continue
            fmt = (getattr(im, "format", None) or "png").lower()
            ext = "." + fmt.lstrip(".")
            imgs.append((ext, blob))
        if imgs:
            by_sheet[ws.title] = imgs
    return by_sheet


def _xlsx_image_bytes(im) -> Optional[bytes]:
    """openpyxl 이미지 객체에서 바이트를 안전하게 추출(버전별 차이 흡수)."""
    # 우선 _data() (로드된 워크북) → ref(BytesIO/경로) 순으로 시도.
    data_fn = getattr(im, "_data", None)
    if callable(data_fn):
        try:
            d = data_fn()
            if d:
                return d
        except Exception:  # noqa: BLE001
            pass
    ref = getattr(im, "ref", None)
    if ref is None:
        return None
    try:
        if isinstance(ref, (bytes, bytearray)):
            return bytes(ref)
        if isinstance(ref, io.BytesIO):
            return ref.getvalue()
        # 경로 문자열/Path
        return Path(str(ref)).read_bytes()
    except Exception:  # noqa: BLE001
        return None


def _inject_xlsx(
    markdown: str,
    path: Path,
    *,
    vision_fn: Callable[..., str],
    provider: Optional[str],
    prompt: str,
) -> str:
    """xlsx markdown 의 각 `## <시트명>` 섹션 뒤에 그 시트 이미지 해설 삽입.

    시트 헤딩을 못 찾으면 말미 모아삽입 fallback.
    """
    by_sheet = _collect_xlsx_images_by_sheet(path)
    if not by_sheet:
        return markdown

    lines = markdown.splitlines()
    # 각 시트 헤딩 라인 위치 찾기: "## <title>" 정확 일치(markitdown xlsx 형식).
    heading_idx: dict[str, int] = {}
    for i, ln in enumerate(lines):
        s = ln.strip()
        if s.startswith("## "):
            title = s[3:].strip()
            if title in by_sheet and title not in heading_idx:
                heading_idx[title] = i

    counter = [0]
    inserts: list[tuple[int, str]] = []
    matched: set[str] = set()
    # 헤딩 위치 정렬: 각 시트 섹션의 끝(다음 헤딩 직전/EOF)에 삽입.
    sorted_headings = sorted(heading_idx.items(), key=lambda kv: kv[1])
    for hi, (title, line_idx) in enumerate(sorted_headings):
        # 섹션 끝 = 다음 "## " 헤딩 직전, 없으면 EOF.
        end = len(lines)
        for j in range(line_idx + 1, len(lines)):
            if lines[j].strip().startswith("## "):
                end = j
                break
        blocks = _describe_many(
            by_sheet[title], label=f"Sheet {title}", counter=counter,
            vision_fn=vision_fn, provider=provider, prompt=prompt,
        )
        if blocks:
            inserts.append((end, "\n\n" + "\n\n".join(blocks)))
        matched.add(title)

    # 헤딩을 못 찾은 시트의 이미지는 말미 모아삽입.
    leftover: list[tuple[str, bytes, str]] = []
    for title in by_sheet:
        if title not in matched:
            for ext, blob in by_sheet[title]:
                leftover.append((ext, blob, f"Sheet {title}"))

    if inserts:
        for pos, text in sorted(inserts, key=lambda t: t[0], reverse=True):
            lines.insert(pos, text)
        markdown = "\n".join(lines)

    if leftover:
        markdown = _append_trailing(
            markdown, leftover, counter=counter,
            vision_fn=vision_fn, provider=provider, prompt=prompt,
        )
    return markdown


# ──────────────────────────────────────────────────────────────────────────────
# docx / hwpx — 말미 모아삽입(근사)
# ──────────────────────────────────────────────────────────────────────────────

def _inject_trailing_from_zip(
    markdown: str,
    path: Path,
    fmt: str,
    *,
    vision_fn: Callable[..., str],
    provider: Optional[str],
    prompt: str,
) -> str:
    """ZIP media 의 이미지를 본문 끝 `## 이미지 해설` 섹션에 순서대로 모아 삽입.

    docx/hwpx 처럼 markitdown 본문에 이미지 출현 위치 마커가 없는(또는 매핑 불가한)
    포맷용 근사 방식. 이미지 0개면 markdown 그대로.
    """
    media = extract_zip_media(path, fmt)
    if not media:
        return markdown
    items: list[tuple[str, bytes, str]] = []
    for name, blob in media:
        ext = Path(name).suffix or ".png"
        items.append((ext, blob, Path(name).name))
    return _append_trailing(markdown, items, vision_fn=vision_fn, provider=provider, prompt=prompt)


# ──────────────────────────────────────────────────────────────────────────────
# 공통: 해설 생성 + 삽입 헬퍼
# ──────────────────────────────────────────────────────────────────────────────

def _describe_many(
    images: list[tuple[str, bytes]],
    *,
    label: str,
    counter: list[int],
    vision_fn: Callable[..., str],
    provider: Optional[str],
    prompt: str,
) -> list[str]:
    """이미지 리스트를 각각 vision 해설하여 통일 형식 블록 리스트로 반환.

    counter[0] 를 문서 단위 이미지 번호로 증가시킨다(실패 이미지는 번호도 소비하지 않음).
    실패 이미지는 skip(블록 미생성).
    """
    blocks: list[str] = []
    for ext, blob in images:
        desc = _describe_image_bytes(
            blob, ext, vision_fn=vision_fn, provider=provider, prompt=prompt
        )
        if desc is None:
            continue
        counter[0] += 1
        blocks.append(_format_block(counter[0], label, desc))
    return blocks


def _append_trailing(
    markdown: str,
    items: list[tuple[str, bytes, str]],
    *,
    vision_fn: Callable[..., str],
    provider: Optional[str],
    prompt: str,
    counter: Optional[list[int]] = None,
) -> str:
    """items(이미지 ext/blob/label)를 본문 끝 `## 이미지 해설` 섹션에 모아 삽입.

    성공 블록이 0개면 섹션을 추가하지 않고 markdown 그대로 반환(빈 섹션 방지).
    """
    if counter is None:
        counter = [0]
    blocks: list[str] = []
    for ext, blob, label in items:
        desc = _describe_image_bytes(
            blob, ext, vision_fn=vision_fn, provider=provider, prompt=prompt
        )
        if desc is None:
            continue
        counter[0] += 1
        blocks.append(_format_block(counter[0], label, desc))
    if not blocks:
        return markdown
    section = _TRAILING_SECTION_TITLE + "\n\n" + "\n\n".join(blocks)
    base = markdown.rstrip()
    if base:
        return base + "\n\n" + section + "\n"
    return section + "\n"


# ──────────────────────────────────────────────────────────────────────────────
# 공개 진입점 — markdown 텍스트 + 원본 경로 → 이미지 해설 삽입된 markdown
# ──────────────────────────────────────────────────────────────────────────────

def augment_with_vision(
    markdown: str,
    src_path: Path | str,
    fmt: str,
    *,
    provider: Optional[str] = None,
    vision_fn: Optional[Callable[..., str]] = None,
    prompt: Optional[str] = None,
) -> str:
    """이미 변환된 markdown 에 src 문서의 이미지 vision 해설을 삽입해 반환한다.

    포맷별 위치 매핑:
      - pptx: 슬라이드 섹션 뒤(정밀)
      - xlsx: 시트 섹션 뒤
      - docx/hwpx: 본문 끝 `## 이미지 해설`(근사)

    어떤 단계든 실패는 safe degrade — 입력 markdown 을 그대로 반환(파이프라인 비중단).

    Args:
        markdown: markitdown 등으로 만든 1차 변환 결과.
        src_path: 원본 office/hwpx 파일 경로(이미지 추출 소스).
        fmt: "docx"|"pptx"|"xlsx"|"hwpx" (소문자).
        provider: vision provider override(None → 모듈 기본 EXTRACT_PROVIDER).
        vision_fn: 이미지→해설 콜러블 override(테스트 mock 용). None 이면
                   fmdw.ollama_extractor.extract_image 를 lazy 로 사용.
        prompt: vision 프롬프트 override(None → FMDW_DESCRIBE_LANG 에 따라 기본
                _IMAGE_VISION_PROMPT_EN(기본)/_IMAGE_VISION_PROMPT_KO 선택).

    Returns:
        이미지 해설이 삽입된 markdown(이미지 0개/전부 실패면 입력 그대로).
    """
    if not markdown:
        markdown = ""
    fmt = (fmt or "").lower()
    if prompt is not None:
        use_prompt = prompt
    else:
        use_prompt = (
            _IMAGE_VISION_PROMPT_KO if _describe_lang() == "ko" else _IMAGE_VISION_PROMPT_EN
        )
    src = Path(src_path)

    try:
        fn = vision_fn if vision_fn is not None else _resolve_vision_fn()
    except Exception as e:  # noqa: BLE001 — vision 엔진 로드 실패 시 텍스트 결과 보존.
        _log.warning("vision 엔진 로드 실패 → 이미지 해설 생략, 텍스트 결과 유지: %s", e)
        return markdown

    try:
        if fmt == "pptx":
            return _inject_pptx(markdown, src, vision_fn=fn, provider=provider, prompt=use_prompt)
        if fmt == "xlsx":
            return _inject_xlsx(markdown, src, vision_fn=fn, provider=provider, prompt=use_prompt)
        if fmt in ("docx", "hwpx"):
            return _inject_trailing_from_zip(
                markdown, src, fmt, vision_fn=fn, provider=provider, prompt=use_prompt
            )
    except Exception as e:  # noqa: BLE001 — 어떤 경로 실패든 텍스트 결과 보존(safe degrade).
        _log.warning("%s 이미지 해설 삽입 실패 → 텍스트 결과 유지: %s", fmt, e)
        return markdown

    # 알 수 없는 포맷 → 텍스트 그대로(무해).
    return markdown
